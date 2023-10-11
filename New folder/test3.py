from pptx import Presentation
from pptx.util import Inches
import win32com.client as win32


# Load the PowerPoint presentation
ppt_app = win32.Dispatch("PowerPoint.Application")

presentation = ppt_app.Presentations.Open("C:\Users\kaosh\Desktop\New folder\power1.pptx")

# Get a specific slide (e.g., slide number 2)
slide = presentation.slides[1]  # Python uses 0-based indexing for lists

# Add a text box to the slide
left = Inches(1)
top = Inches(1)
width = Inches(8)
height = Inches(1.5)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame

# Add text to the text box
text = "Hello, PowerPoint!"
text_frame.text = text

# Save the presentation with the modified content
presentation.save("path/to/modified_presentation.pptx")
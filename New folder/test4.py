from pptx import Presentation
from pptx.util import Pt,Cm

# 打開已存在ppt
ppt = Presentation('power1.pptx')

# 設置添加到當前ppt哪一頁
n_page = 0
singleLineContent = "我是單行內容"
multiLineContent = \
"""我是多行內容1
我是多行內容2
我是多行內容3
"""

# 獲取需要添加文字的頁面對象
slide = ppt.slides[n_page]

# 添加單行內容

# 設置添加文字框的位置以及大小
left, top, width, height = Cm(16.9), Cm(1), Cm(12), Cm(1.2)
# 添加文字段落
new_paragraph1 = slide.shapes.add_textbox(left=left, top=top, width=width, height=height).text_frame
# 設置段落內容
new_paragraph1.paragraphs[0].text = singleLineContent
# 設置文字大小
new_paragraph1.paragraphs[0].font.size = Pt(15)


# 添加多行

# 設置添加文字框的位置以及大小
left, top, width, height = Cm(16.9), Cm(3), Cm(12), Cm(3.6)
# 添加文字段落
new_paragraph2 = slide.shapes.add_textbox(left=left, top=top, width=width, height=height).text_frame
# 設置段落內容
new_paragraph2.paragraphs[0].text = multiLineContent
# 設置文字大小
new_paragraph2.paragraphs[0].font.size = Pt(15)


# 保存ppt
#ppt.save('4.1 添加文字.pptx')